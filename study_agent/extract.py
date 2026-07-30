"""教科書・授業資料からのテキスト抽出 (PDF / docx / pptx / txt / md)。"""
from pathlib import Path


def extract_text(path: str) -> str:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"ファイルが見つかりません: {path}")
    ext = p.suffix.lower()

    if ext == ".pdf":
        return _from_pdf(p)
    if ext == ".docx":
        return _from_docx(p)
    if ext == ".pptx":
        return _from_pptx(p)
    if ext in (".txt", ".md", ".markdown"):
        return p.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"未対応の形式です: {ext} (pdf/docx/pptx/txt/md に対応)")


def _from_pdf(p: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        pages.append(f"[p.{i}]\n{text}")
    return "\n\n".join(pages)


def _from_docx(p: Path) -> str:
    import docx

    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs if para.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _from_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
        block = f"[Slide {i}]\n" + "\n".join(t for t in texts if t.strip())
        if notes.strip():
            block += f"\n(Speaker notes: {notes})"
        slides.append(block)
    return "\n\n".join(slides)


def chunk_text(text: str, max_chars: int = 12000) -> list[str]:
    """段落境界を優先して長文を分割する。"""
    if len(text) <= max_chars:
        return [text]
    chunks, current = [], []
    size = 0
    for para in text.split("\n\n"):
        if size + len(para) > max_chars and current:
            chunks.append("\n\n".join(current))
            current, size = [], 0
        current.append(para)
        size += len(para) + 2
    if current:
        chunks.append("\n\n".join(current))
    return chunks
